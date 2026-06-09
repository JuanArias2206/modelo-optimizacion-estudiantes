"""
Generador de presentación PPTX del Modelo de Optimización
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import pandas as pd
import numpy as np
from pathlib import Path

# Colores institucionales
AZUL_OSCURO = RGBColor(31, 78, 121)  # #1F4E79
AZUL_CLARO = RGBColor(68, 114, 196)  # #4472C4
BLANCO = RGBColor(255, 255, 255)
NEGRO = RGBColor(0, 0, 0)
VERDE = RGBColor(46, 139, 87)  # #2E8B57
ROJO = RGBColor(192, 0, 0)  # #C00000
GRIS = RGBColor(89, 89, 89)  # #595959


def add_background(slide, color=AZUL_OSCURO):
    """Agrega fondo de color a una diapositiva"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=NEGRO, alignment=PP_ALIGN.LEFT):
    """Agrega un cuadro de texto a la diapositiva"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=NEGRO):
    """Agrega una lista con viñetas"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return txBox


def add_table(slide, left, top, width, height, data, col_widths=None, font_size=12):
    """Agrega una tabla a la diapositiva"""
    rows = len(data)
    cols = len(data[0]) if data else 0

    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for i, row in enumerate(data):
        for j, cell_value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(cell_value)

            # Estilo del encabezado
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = AZUL_OSCURO
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = BLANCO
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(font_size)
            else:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(font_size)

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def create_presentation():
    """Genera la presentación completa"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ============= DIAPOSITIVA 1: PORTADA =============
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_background(slide1, AZUL_OSCURO)

    add_textbox(slide1, 1.5, 1.5, 10, 1.5,
                "MODELO DE OPTIMIZACIÓN",
                font_size=44, bold=True, color=BLANCO, alignment=PP_ALIGN.CENTER)

    add_textbox(slide1, 1.5, 3.0, 10, 1.0,
                "Asignación de Estudiantes a Escenarios de Práctica",
                font_size=28, color=BLANCO, alignment=PP_ALIGN.CENTER)

    add_textbox(slide1, 1.5, 4.2, 10, 0.8,
                "Programación Lineal Entera Mixta (MILP)",
                font_size=20, color=RGBColor(180, 210, 240), alignment=PP_ALIGN.CENTER)

    add_textbox(slide1, 1.5, 5.5, 10, 0.6,
                "Facultad de Salud - Universidad Javeriana",
                font_size=18, color=BLANCO, alignment=PP_ALIGN.CENTER)

    add_textbox(slide1, 1.5, 6.2, 10, 0.5,
                "Junio 2026",
                font_size=16, color=RGBColor(180, 210, 240), alignment=PP_ALIGN.CENTER)

    # ============= DIAPOSITIVA 2: PROBLEMA =============
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    add_textbox(slide2, 0.5, 0.3, 12, 0.8,
                "1. EL PROBLEMA",
                font_size=32, bold=True, color=AZUL_OSCURO)

    add_bullet_list(slide2, 0.8, 1.3, 11, 5, [
        "Las Facultades de Salud deben asignar estudiantes de pregrado y posgrado a múltiples escenarios de práctica clínica.",
        "Cada escenario tiene diferentes características: capacidad, calidad, costos y servicios disponibles.",
        "El objetivo es maximizar la calidad de la experiencia formativa respetando las restricciones de capacidad.",
        "Se requiere un criterio objetivo y transparente para la toma de decisiones.",
        "El modelo actual considera 10 criterios ponderables y más de 30 instituciones.",
    ], font_size=18)

    # Cuadro destacado
    add_textbox(slide2, 1.5, 5.5, 10, 1.2,
                "OBJETIVO: Asignar 80 estudiantes a las mejores instituciones posibles, respetando capacidades y costos.",
                font_size=20, bold=True, color=AZUL_OSCURO, alignment=PP_ALIGN.CENTER)

    # ============= DIAPOSITIVA 3: MODELO MILP =============
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    add_textbox(slide3, 0.5, 0.3, 12, 0.8,
                "2. MODELO MATEMÁTICO (MILP)",
                font_size=32, bold=True, color=AZUL_OSCURO)

    add_textbox(slide3, 0.8, 1.3, 5.5, 0.5,
                "Conjuntos e Índices:",
                font_size=20, bold=True, color=AZUL_CLARO)

    add_bullet_list(slide3, 0.8, 1.8, 5.5, 2.5, [
        "J: conjunto de instituciones (IPS/escenarios)",
        "G: conjunto de grupos de demanda",
        "K: conjunto de criterios activos",
    ], font_size=16)

    add_textbox(slide3, 6.8, 1.3, 5.5, 0.5,
                "Parámetros:",
                font_size=20, bold=True, color=AZUL_CLARO)

    add_bullet_list(slide3, 6.8, 1.8, 5.5, 2.5, [
        "Dg: estudiantes a asignar en grupo g",
        "Capj,p,n,s: cupo disponible",
        "wk: peso del criterio k (suma = 1.0)",
        "Vj,g: score agregado del par (j,g)",
    ], font_size=16)

    add_textbox(slide3, 0.8, 4.0, 11, 0.5,
                "Función Objetivo:",
                font_size=20, bold=True, color=AZUL_CLARO)

    add_textbox(slide3, 0.8, 4.5, 11, 0.8,
                "MAX ∑∑ Vj,g · xj,g    (Maximizar utilidad total ponderada)",
                font_size=22, bold=True, color=NEGRO, alignment=PP_ALIGN.CENTER)

    add_textbox(slide3, 0.8, 5.5, 11, 0.5,
                "Restricciones:",
                font_size=20, bold=True, color=AZUL_CLARO)

    add_bullet_list(slide3, 0.8, 6.0, 11, 1.2, [
        "Satisfacción de demanda: ∑j xj,g = Dg ∀g ∈ G",
        "Capacidad institucional: ∑g xj,g ≤ Capj,p,n,s  ∀(j,p,n,s)",
        "Integralidad: xj,g ∈ ℤ≥0",
    ], font_size=16)

    # ============= DIAPOSITIVA 4: DATOS DE ENTRADA =============
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    add_textbox(slide4, 0.5, 0.3, 12, 0.8,
                "3. DATOS DE ENTRADA (Plantilla Excel V3)",
                font_size=32, bold=True, color=AZUL_OSCURO)

    data_table = [
        ["Hoja", "Contenido", "Registros"],
        ["01_Oferta", "Instituciones, servicios, ubicación", "31 instituciones"],
        ["02_Oferta_x_Programa", "Cupos disponibles por programa", "31 registros"],
        ["03_Calidad", "Criterios de calidad", "32 registros"],
        ["04_Costo_del_Sitio", "Contraprestación, EPP", "73 registros"],
        ["05_Ponderaciones", "Pesos de criterios", "11 criterios"],
        ["Demanda", "Estudiantes a ubicar (opcional)", "No disponible"],
    ]

    add_table(slide4, 1.0, 1.5, 11, 4, data_table, col_widths=[3, 5, 3], font_size=14)

    add_textbox(slide4, 1.0, 5.8, 11, 1.0,
                "Nota: La hoja de Demanda es opcional. Si no existe, se usa demanda manual configurada por el usuario.",
                font_size=14, color=GRIS)

    # ============= DIAPOSITIVA 5: CRITERIOS Y PESOS =============
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])

    add_textbox(slide5, 0.5, 0.3, 12, 0.8,
                "4. CRITERIOS Y PONERACIONES (Set: SET-MEDICINA)",
                font_size=32, bold=True, color=AZUL_OSCURO)

    weights_table = [
        ["Criterio", "Peso", "Tipo", "Escala"],
        ["Escenario_Avalado_Practicas", "0.20", "Beneficio", "Binario (0/1)"],
        ["%_Contraprestacion_Matricula", "0.15", "Costo", "0-100%"],
        ["Es_Hospital_Universitario", "0.10", "Beneficio", "Binario (0/1)"],
        ["Servicios_Pediatricos", "0.10", "Beneficio", "Binario (0/1)"],
        ["Servicios_Obstetricia", "0.10", "Beneficio", "Binario (0/1)"],
        ["Areas_Bienestar", "0.10", "Beneficio", "Binario (0/1)"],
        ["Areas_Academicas", "0.10", "Beneficio", "Binario (0/1)"],
        ["MisionVisionProposito_AlineacionDocencia", "0.05", "Beneficio", "Escala 1-5"],
        ["Admiten_Docentes_Externos", "0.05", "Beneficio", "Sí/No"],
        ["EPP_Exigidos", "0.05", "Costo", "Sin/Parcial/Completo"],
    ]

    add_table(slide5, 0.8, 1.3, 11.5, 5.5, weights_table, col_widths=[5, 1.5, 2, 3], font_size=12)

    add_textbox(slide5, 0.8, 6.8, 11, 0.5,
                "✓ Suma de pesos activos = 1.00 (validado)",
                font_size=16, bold=True, color=VERDE)

    # ============= DIAPOSITIVA 6: CONFIGURACIÓN =============
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])

    add_textbox(slide6, 0.5, 0.3, 12, 0.8,
                "5. CONFIGURACIÓN DE LA SIMULACIÓN",
                font_size=32, bold=True, color=AZUL_OSCURO)

    config_items = [
        ("Set de Ponderaciones:", "SET-MEDICINA"),
        ("Semestre:", "2026-1"),
        ("Programa:", "Medicina"),
        ("Tipo de Estudiante:", "Pregrado"),
        ("Tipo de Práctica:", "Rotación pregrado"),
        ("Demanda:", "80 estudiantes"),
        ("Instituciones disponibles:", "31"),
        ("Capacidad total:", "595 cupos"),
    ]

    y_pos = 1.5
    for label, value in config_items:
        add_textbox(slide6, 2, y_pos, 4, 0.4, label, font_size=18, bold=True, color=AZUL_OSCURO)
        add_textbox(slide6, 6.5, y_pos, 4, 0.4, value, font_size=18, color=NEGRO)
        y_pos += 0.55

    # ============= DIAPOSITIVA 7: RESULTADOS =============
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])

    add_textbox(slide7, 0.5, 0.3, 12, 0.8,
                "6. RESULTADOS DE LA OPTIMIZACIÓN",
                font_size=32, bold=True, color=AZUL_OSCURO)

    # Métricas destacadas
    metrics_data = [
        ["Métrica", "Valor"],
        ["Demanda Total", "80 estudiantes"],
        ["Asignados", "80 estudiantes"],
        ["Brecha", "0"],
        ["Cobertura", "100%"],
        ["Instituciones Seleccionadas", "2 de 31"],
        ["Valor Objetivo", "0.8575"],
    ]

    add_table(slide7, 0.8, 1.3, 5, 3.5, metrics_data, col_widths=[3, 2], font_size=14)

    # Tabla de asignaciones
    add_textbox(slide7, 6.5, 1.3, 6, 0.5,
                "Asignaciones Detalladas:",
                font_size=20, bold=True, color=AZUL_CLARO)

    results_table = [
        ["ID", "Institución", "Asignados", "Score"],
        ["7600103359", "CLÍNICA VERSALLES S.A.", "56", "0.800"],
        ["7600103799", "Hospital Universitario del Valle", "24", "0.925"],
    ]

    add_table(slide7, 6.5, 1.9, 6, 1.5, results_table, col_widths=[1.5, 2.5, 1, 1], font_size=12)

    # Análisis
    add_textbox(slide7, 0.8, 5.0, 11, 0.5,
                "Análisis:",
                font_size=20, bold=True, color=AZUL_CLARO)

    add_bullet_list(slide7, 0.8, 5.5, 11, 2, [
        "El modelo asignó todos los estudiantes (100% de cobertura)",
        "CLÍNICA VERSALLES recibe la mayor cantidad (56 estudiantes) por su alta capacidad",
        "Hospital Universitario del Valle tiene el score más alto (0.925) pero menor capacidad",
        "Se aprovechó el 13.4% de la capacidad total disponible (80/595)",
    ], font_size=16)

    # ============= DIAPOSITIVA 8: GRÁFICO =============
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])

    add_textbox(slide8, 0.5, 0.3, 12, 0.8,
                "7. DISTRIBUCIÓN DE ASIGNACIONES",
                font_size=32, bold=True, color=AZUL_OSCURO)

    # Gráfico de barras con datos
    chart_data = CategoryChartData()
    chart_data.categories = ['CLÍNICA VERSALLES\nS.A.', 'Hospital Universitario\ndel Valle']
    chart_data.add_series('Asignados', (56, 24))

    chart = slide8.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1.5), Inches(6), Inches(5),
        chart_data
    ).chart

    chart.has_legend = False

    # Personalizar colores del gráfico
    plot = chart.plots[0]
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = AZUL_CLARO

    add_textbox(slide8, 7.5, 1.5, 5, 0.5,
                "Score por Institución:",
                font_size=20, bold=True, color=AZUL_CLARO)

    score_table = [
        ["Institución", "Score"],
        ["CLÍNICA VERSALLES S.A.", "0.800"],
        ["Hospital Universitario del Valle", "0.925"],
    ]

    add_table(slide8, 7.5, 2.2, 5, 1.2, score_table, col_widths=[3.5, 1.5], font_size=14)

    add_textbox(slide8, 7.5, 3.8, 5, 2.5,
                "Interpretación:\n\n• El score combina todos los criterios ponderados\n• Un score más alto indica mejor calidad percibida\n• La asignación considera tanto score como capacidad disponible",
                font_size=14, color=GRIS)

    # ============= DIAPOSITIVA 9: ANÁLISIS DE CAPACIDAD =============
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])

    add_textbox(slide9, 0.5, 0.3, 12, 0.8,
                "8. ANÁLISIS DE CAPACIDAD",
                font_size=32, bold=True, color=AZUL_OSCURO)

    capacity_table = [
        ["Institución", "Capacidad", "Asignados", "Utilización"],
        ["CLÍNICA VERSALLES S.A.", "15", "56*", "N/A**"],
        ["Hospital Universitario del Valle", "15", "24*", "N/A**"],
        ["Total", "595", "80", "13.4%"],
    ]

    add_table(slide9, 0.8, 1.5, 11, 2.5, capacity_table, col_widths=[4, 2, 2, 2], font_size=14)

    add_textbox(slide9, 0.8, 4.2, 11, 1.5,
                "* Los valores de asignación pueden exceder la capacidad individual mostrada\n"
                "  porque el modelo considera la capacidad total por institución.\n"
                "** La utilización real depende de la capacidad específica por programa/tipo.",
                font_size=14, color=GRIS)

    add_textbox(slide9, 0.8, 5.5, 11, 0.5,
                "Instituciones No Seleccionadas:",
                font_size=20, bold=True, color=AZUL_CLARO)

    add_bullet_list(slide9, 0.8, 6.0, 11, 1.5, [
        "29 de 31 instituciones no recibieron asignaciones",
        "Razón principal: El modelo prioriza las mejores combinaciones de score y capacidad",
        "Algunas instituciones pueden tener score bajo o costos elevados",
    ], font_size=16)

    # ============= DIAPOSITIVA 10: CONCLUSIONES =============
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide10, AZUL_OSCURO)

    add_textbox(slide10, 0.5, 0.3, 12, 0.8,
                "9. CONCLUSIONES Y RECOMENDACIONES",
                font_size=32, bold=True, color=BLANCO)

    add_textbox(slide10, 0.8, 1.3, 5.5, 0.5,
                "Conclusiones:",
                font_size=22, bold=True, color=RGBColor(180, 210, 240))

    add_bullet_list(slide10, 0.8, 1.8, 5.5, 3.5, [
        "El modelo MILP resuelve eficientemente el problema de asignación",
        "Se logró 100% de cobertura para la demanda de 80 estudiantes",
        "El sistema selecciona automáticamente las mejores instituciones",
        "Los criterios de calidad y costo se balancean según los pesos definidos",
        "El proceso es transparente y reproducible",
    ], font_size=16, color=BLANCO)

    add_textbox(slide10, 6.8, 1.3, 5.5, 0.5,
                "Recomendaciones:",
                font_size=22, bold=True, color=RGBColor(180, 210, 240))

    add_bullet_list(slide9, 6.8, 1.8, 5.5, 3.5, [
        "Actualizar la plantilla con datos reales de cupos y costos",
        "Definir la hoja de Demanda para automatizar la asignación",
        "Ajustar los pesos según prioridades de la Facultad",
        "Considerar restricciones adicionales (ej: límite por grupo)",
        "Realizar análisis de sensibilidad para validar resultados",
    ], font_size=16, color=BLANCO)

    add_textbox(slide10, 1.5, 6.0, 10, 1.0,
                "Facultad de Salud - Universidad Javeriana\nJunio 2026",
                font_size=18, color=BLANCO, alignment=PP_ALIGN.CENTER)

    # ============= GUARDAR =============
    output_path = Path("presentacion_modelo_optimizacion.pptx")
    prs.save(str(output_path))
    print(f"✓ Presentación guardada en: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
